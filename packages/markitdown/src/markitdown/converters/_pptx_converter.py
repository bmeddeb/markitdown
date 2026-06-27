import sys
import base64
import os
import io
import re
import html

from typing import BinaryIO, Any
from operator import attrgetter

from ._html_converter import HtmlConverter
from ._llm_caption import llm_caption
from .._base_converter import DocumentConverter, DocumentConverterResult
from .._stream_info import StreamInfo
from .._exceptions import MissingDependencyException, MISSING_DEPENDENCY_MESSAGE

# Try loading optional (but in this case, required) dependencies
# Save reporting of any exceptions for later
_dependency_exc_info = None
try:
    import pptx
except ImportError:
    # Preserve the error and stack trace for later
    _dependency_exc_info = sys.exc_info()


ACCEPTED_MIME_TYPE_PREFIXES = [
    "application/vnd.openxmlformats-officedocument.presentationml",
]

ACCEPTED_FILE_EXTENSIONS = [".pptx"]


class PptxConverter(DocumentConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def __init__(self):
        super().__init__()
        self._html_converter = HtmlConverter()

    def accepts(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> bool:
        mimetype = (stream_info.mimetype or "").lower()
        extension = (stream_info.extension or "").lower()

        if extension in ACCEPTED_FILE_EXTENSIONS:
            return True

        for prefix in ACCEPTED_MIME_TYPE_PREFIXES:
            if mimetype.startswith(prefix):
                return True

        return False

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> DocumentConverterResult:
        # Check the dependencies
        if _dependency_exc_info is not None:
            raise MissingDependencyException(
                MISSING_DEPENDENCY_MESSAGE.format(
                    converter=type(self).__name__,
                    extension=".pptx",
                    feature="pptx",
                )
            ) from _dependency_exc_info[
                1
            ].with_traceback(  # type: ignore[union-attr]
                _dependency_exc_info[2]
            )

        # Perform the conversion
        presentation = pptx.Presentation(file_stream)
        md_content = ""
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            def append_block(block):
                nonlocal md_content
                block = block.strip()
                if not block:
                    return
                if md_content:
                    md_content = md_content.rstrip() + "\n\n"
                md_content += block + "\n\n"

            sorted_shapes = sorted(
                slide.shapes,
                key=lambda x: (
                    float("-inf") if not x.top else x.top,
                    float("-inf") if not x.left else x.left,
                ),
            )
            title = self._find_title_shape(slide, sorted_shapes)
            title_text = (
                self._clean_heading_text(title.text) if title is not None else ""
            )
            if title_text:
                append_block(f"## Slide {slide_num}: {title_text}")
            else:
                append_block(f"## Slide {slide_num}")

            def get_shape_content(shape, **kwargs):
                nonlocal md_content
                if shape == title:
                    return

                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069

                    llm_description = ""
                    alt_text = ""

                    # Potentially generate a description using an LLM
                    llm_client = kwargs.get("llm_client")
                    llm_model = kwargs.get("llm_model")
                    if llm_client is not None and llm_model is not None:
                        # Prepare a file_stream and stream_info for the image data
                        image_filename = shape.image.filename
                        image_extension = None
                        if image_filename:
                            image_extension = os.path.splitext(image_filename)[1]
                        image_stream_info = StreamInfo(
                            mimetype=shape.image.content_type,
                            extension=image_extension,
                            filename=image_filename,
                        )

                        image_stream = io.BytesIO(shape.image.blob)

                        # Caption the image
                        try:
                            llm_description = llm_caption(
                                image_stream,
                                image_stream_info,
                                client=llm_client,
                                model=llm_model,
                                prompt=kwargs.get("llm_prompt"),
                            )
                        except Exception:
                            # Unable to generate a description
                            pass

                    # Also grab any description embedded in the deck
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        # Unable to get alt text
                        pass

                    # Prepare the alt, escaping any special characters
                    alt_text = "\n".join([llm_description, alt_text]) or shape.name
                    alt_text = re.sub(r"[\r\n\[\]]", " ", alt_text)
                    alt_text = re.sub(r"\s+", " ", alt_text).strip()

                    # If keep_data_uris is True, use base64 encoding for images
                    if kwargs.get("keep_data_uris", False):
                        blob = shape.image.blob
                        content_type = shape.image.content_type or "image/png"
                        b64_string = base64.b64encode(blob).decode("utf-8")
                        append_block(
                            f"![{alt_text}](data:{content_type};base64,{b64_string})"
                        )
                    else:
                        # A placeholder name
                        filename = re.sub(r"\W", "", shape.name) + ".jpg"
                        append_block("![" + alt_text + "](" + filename + ")")

                # Tables
                if self._is_table(shape):
                    append_block(self._convert_table_to_markdown(shape.table, **kwargs))

                # Charts
                if shape.has_chart:
                    append_block(self._convert_chart_to_markdown(shape.chart))

                # Text areas
                elif shape.has_text_frame:
                    append_block(self._convert_text_frame_to_markdown(shape.text_frame))

                # Group Shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                    sorted_shapes = sorted(
                        shape.shapes,
                        key=lambda x: (
                            float("-inf") if not x.top else x.top,
                            float("-inf") if not x.left else x.left,
                        ),
                    )
                    for subshape in sorted_shapes:
                        get_shape_content(subshape, **kwargs)

            for shape in sorted_shapes:
                get_shape_content(shape, **kwargs)

            md_content = md_content.strip()

            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    notes_markdown = self._convert_text_frame_to_markdown(notes_frame)
                    if notes_markdown:
                        md_content = md_content.rstrip() + "\n\n"
                        md_content += f"### Notes\n\n{notes_markdown}"
                md_content = md_content.strip()

        return DocumentConverterResult(markdown=md_content.strip())

    def _find_title_shape(self, slide, sorted_shapes):
        title = slide.shapes.title
        if title is not None and self._clean_text(title.text):
            return title

        for shape in sorted_shapes:
            if self._is_title_placeholder(shape) and self._clean_text(shape.text):
                return shape

        for shape in sorted_shapes:
            if self._looks_like_title_shape(shape):
                return shape

        return None

    def _is_title_placeholder(self, shape):
        if not getattr(shape, "is_placeholder", False):
            return False

        try:
            placeholder_type = shape.placeholder_format.type
        except Exception:
            return False

        return placeholder_type in (
            pptx.enum.shapes.PP_PLACEHOLDER.CENTER_TITLE,
            pptx.enum.shapes.PP_PLACEHOLDER.TITLE,
            pptx.enum.shapes.PP_PLACEHOLDER.VERTICAL_TITLE,
        )

    def _looks_like_title_shape(self, shape):
        if not getattr(shape, "has_text_frame", False):
            return False

        text = self._clean_text(shape.text)
        if not text or len(text) > 160:
            return False

        non_empty_paragraphs = [
            paragraph
            for paragraph in shape.text_frame.paragraphs
            if self._clean_text(paragraph.text)
        ]
        if len(non_empty_paragraphs) != 1:
            return False

        if self._is_non_content_placeholder(shape):
            return False

        return True

    def _is_non_content_placeholder(self, shape):
        if not getattr(shape, "is_placeholder", False):
            return False

        try:
            placeholder_type = shape.placeholder_format.type
        except Exception:
            return False

        return placeholder_type in (
            pptx.enum.shapes.PP_PLACEHOLDER.DATE,
            pptx.enum.shapes.PP_PLACEHOLDER.FOOTER,
            pptx.enum.shapes.PP_PLACEHOLDER.HEADER,
            pptx.enum.shapes.PP_PLACEHOLDER.SLIDE_NUMBER,
        )

    def _convert_text_frame_to_markdown(self, text_frame):
        blocks = []
        current_list = []

        def flush_list():
            if current_list:
                blocks.append("\n".join(current_list))
                current_list.clear()

        for paragraph in text_frame.paragraphs:
            text = self._clean_text(paragraph.text)
            if not text:
                flush_list()
                continue

            if self._is_bullet_paragraph(paragraph):
                indent = "  " * max(paragraph.level, 0)
                current_list.append(f"{indent}- {text}")
            else:
                flush_list()
                blocks.append(text)

        flush_list()
        return "\n\n".join(blocks)

    def _is_bullet_paragraph(self, paragraph):
        if paragraph.level > 0:
            return True

        p_pr = getattr(paragraph._p, "pPr", None)
        if p_pr is None:
            return False

        has_bullet = False
        for child in p_pr:
            tag_name = child.tag.rsplit("}", 1)[-1]
            if tag_name == "buNone":
                return False
            if tag_name in ("buAutoNum", "buBlip", "buChar"):
                has_bullet = True

        return has_bullet

    def _clean_text(self, text):
        text = (text or "").replace("\xa0", " ")
        lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
        return "\n".join(line for line in lines if line).strip()

    def _clean_heading_text(self, text):
        return re.sub(r"\s+", " ", self._clean_text(text)).strip()

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_table_to_markdown(self, table, **kwargs):
        # Write the table as HTML, then convert it to Markdown
        html_table = "<html><body><table>"
        first_row = True
        for row in table.rows:
            html_table += "<tr>"
            for cell in row.cells:
                if first_row:
                    html_table += "<th>" + html.escape(cell.text) + "</th>"
                else:
                    html_table += "<td>" + html.escape(cell.text) + "</td>"
            html_table += "</tr>"
            first_row = False
        html_table += "</table></body></html>"

        return (
            self._html_converter.convert_string(html_table, **kwargs).markdown.strip()
            + "\n"
        )

    def _convert_chart_to_markdown(self, chart):
        try:
            md = "\n\n### Chart"
            if chart.has_title:
                md += f": {chart.chart_title.text_frame.text}"
            md += "\n\n"
            data = []
            category_names = [c.label for c in chart.plots[0].categories]
            series_names = [s.name for s in chart.series]
            data.append(["Category"] + series_names)

            for idx, category in enumerate(category_names):
                row = [category]
                for series in chart.series:
                    row.append(series.values[idx])
                data.append(row)

            markdown_table = []
            for row in data:
                markdown_table.append("| " + " | ".join(map(str, row)) + " |")
            header = markdown_table[0]
            separator = "|" + "|".join(["---"] * len(data[0])) + "|"
            return md + "\n".join([header, separator] + markdown_table[1:])
        except ValueError as e:
            # Handle the specific error for unsupported chart types
            if "unsupported plot type" in str(e):
                return "\n\n[unsupported chart]\n\n"
        except Exception:
            # Catch any other exceptions that might occur
            return "\n\n[unsupported chart]\n\n"
